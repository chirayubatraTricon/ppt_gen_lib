import os
from pptx import Presentation
from typing import List, Dict

from pathlib import Path
from ..state import PptGenerationState
from ..base.node import Node

class InspectTemplateLayout(Node[PptGenerationState]):
    name = "save_uploaded_path"

    async def __call__(self, state: PptGenerationState ):
        try:
            """Return layouts or raise an exception with a helpful message."""
            template_path = state["template_path"]
            layouts = self._list_template_layouts(template_path)
            if not layouts:
                raise RuntimeError("The uploaded template has no slide layouts.")
            state["layouts"] = layouts

        except Exception as e:
            print("error occured in saving file path")

        return state
    
    def _list_template_layouts(self, template_path: str | Path) -> List[Dict]:
        """Return simple metadata for each slide layout in a .pptx template.

        Each returned dict has keys: 'index', 'name', 'placeholders' (list of names).

        Raises FileNotFoundError or ValueError for invalid inputs.
        """
        path = Path(template_path)
        if not path.exists():
            raise FileNotFoundError(f"Template not found: {path}")
        if path.suffix.lower() != ".pptx":
            raise ValueError("Only .pptx files are supported")

        prs = Presentation(str(path))
        result: List[Dict] = []

        for idx, layout in enumerate(prs.slide_layouts):
            ph_names: List[str] = []
            for ph in getattr(layout, "placeholders", []):
                name = getattr(ph, "name", None) or f"placeholder_{getattr(getattr(ph, 'placeholder_format', None), 'idx', 'unk')}"
                ph_names.append(name)

            result.append({
                "index": idx,
                "name": layout.name or f"Layout {idx}",
                "placeholders": ph_names,
            })

        return result